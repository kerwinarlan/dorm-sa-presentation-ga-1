import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    photos_dir = os.path.join(base_dir, "assets", "staff_photos")

    # Theme Colors
    RED = RGBColor(211, 47, 47)        # #D32F2F
    DARK_RED = RGBColor(139, 28, 28)   # #8B1C1C
    DEEP_BG = RGBColor(60, 10, 10)     # #3C0A0A
    DARK_CARD = RGBColor(30, 30, 36)   # #1E1E24
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(18, 18, 18)
    YELLOW = RGBColor(251, 192, 45)

    blank_layout = prs.slide_layouts[6]

    def add_bg(slide, color=DEEP_BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_badge(slide, text, top=Inches(0.5)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), top, Inches(2.8), Inches(0.45))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RED
        shape.line.color.rgb = WHITE
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        return shape

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1, DEEP_BG)

    card1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card1.fill.solid()
    card1.fill.fore_color.rgb = DARK_RED
    card1.line.color.rgb = WHITE
    card1.line.width = Pt(3)

    tf1 = card1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "SAMPAGUITA RESIDENCE HALL"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = YELLOW
    p.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "\nGENERAL ASSEMBLY"
    p2.font.size = Pt(54)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf1.add_paragraph()
    p3.text = "PLAN OF ACTIVITIES"
    p3.font.size = Pt(36)
    p3.font.bold = True
    p3.font.color.rgb = YELLOW
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf1.add_paragraph()
    p4.text = "\nFirst Semester, A.Y. 2026-2027"
    p4.font.size = Pt(20)
    p4.font.color.rgb = WHITE
    p4.alignment = PP_ALIGN.CENTER

    # SLIDE 2: Dorm Managers & Houseparent
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_badge(slide2, "ADMINISTRATION")

    txBox = slide2.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "DORM MANAGERS & HOUSEPARENT"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    admins = [
        ("Ma'am Len", "Dorm Manager", "maam-len.png"),
        ("Ma'am Zay", "OIC / Dorm Manager", "maam-zay.png"),
        ("Ma'am Gen", "Houseparent", "maam-gen.png")
    ]

    for idx, (name, role, img_name) in enumerate(admins):
        left = Inches(1.0 + idx * 3.9)
        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), Inches(3.5), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = YELLOW
        card.line.width = Pt(3)

        img_path = os.path.join(photos_dir, img_name)
        if os.path.exists(img_path):
            slide2.shapes.add_picture(img_path, left + Inches(0.5), Inches(2.2), Inches(2.5), Inches(2.6))

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"\n\n\n\n\n\n\n\n{role.upper()}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"{name}"
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = BLACK
        p2.alignment = PP_ALIGN.CENTER

    # SLIDE 3: Student Assistants
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_badge(slide3, "YOUR SA TEAM")

    txBox = slide3.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "STUDENT ASSISTANTS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sas = [
        ("Yelsah", "yelsah.png"),
        ("Rachel", "rachel.png"),
        ("Jose", "jose.png"),
        ("Ash (GA1 Head)", "ash.png"),
        ("Matt", "matt.png")
    ]
    for idx, (sa, img_name) in enumerate(sas):
        left = Inches(0.8 + idx * 2.4)
        card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), Inches(2.1), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RED if "Ash" in sa else BLACK
        card.line.width = Pt(3)

        img_path = os.path.join(photos_dir, img_name)
        if os.path.exists(img_path):
            slide3.shapes.add_picture(img_path, left + Inches(0.15), Inches(2.2), Inches(1.8), Inches(2.2))

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"\n\n\n\n\n\n\nSA"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"{sa}"
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = BLACK
        p2.alignment = PP_ALIGN.CENTER

    # SLIDE 4: Office & Maintenance Staff
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_badge(slide4, "SUPPORT TEAM")

    txBox = slide4.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "OFFICE & MAINTENANCE STAFF"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    staff_members = [
        ("Ate Rhea", "ate-rhea.png"),
        ("Ate December", "ate-december.png"),
        ("Ate Feny", "ate-feny.png"),
        ("Ate Tere", "ate-tere.png"),
        ("Kuya Bagyo", "kuya-bagyo.png"),
        ("Kuya Dani", "kuya-dani.png"),
        ("Kuya Emong", "kuya-emong.png")
    ]

    for idx, (name, img_name) in enumerate(staff_members):
        left = Inches(0.6 + idx * 1.75)
        card = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.2), Inches(1.6), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BLACK
        card.line.width = Pt(2)

        img_path = os.path.join(photos_dir, img_name)
        if os.path.exists(img_path):
            slide4.shapes.add_picture(img_path, left + Inches(0.1), Inches(2.4), Inches(1.4), Inches(1.8))

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"\n\n\n\n\n\n\n{name}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER

    # SLIDE 5: Security Personnel
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_badge(slide5, "PROTECTORS")

    txBox = slide5.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "SECURITY PERSONNEL"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    guards = [
        ("Kuya Nard", "kuya-nard.png"),
        ("Kuya Julius", "kuya-julius.png"),
        ("Ate Jen", "ate-jen.png"),
        ("Ate Joh", "ate-joh.png")
    ]

    for idx, (name, img_name) in enumerate(guards):
        left = Inches(1.0 + idx * 2.9)
        card = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), Inches(2.6), Inches(4.6))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BLACK
        card.line.width = Pt(3)

        img_path = os.path.join(photos_dir, img_name)
        if os.path.exists(img_path):
            slide5.shapes.add_picture(img_path, left + Inches(0.2), Inches(2.2), Inches(2.2), Inches(2.5))

        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"\n\n\n\n\n\n\n\nSECURITY GUARD"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RED
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"{name}"
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = BLACK
        p2.alignment = PP_ALIGN.CENTER

    # SLIDE 6: Overview / Agenda
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_badge(slide6, "OVERVIEW")

    txBox = slide6.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PLAN OF ACTIVITIES OVERVIEW"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    categories = [
        ("01", "GA1 / AKWE", "Aug 31 - Sep 4", "SPIDER-MAN THEME", RED),
        ("02", "HALLOWEEN WEEK", "Oct 19 - Oct 23", "HORROR THEME", DARK_CARD),
        ("03", "GA2 / YEAR-END", "End of Semester", "CELEBRATION", RED),
        ("04", "ADDITIONAL EVENTS", "Semester-Long", "SPECIAL PROJECTS", DARK_CARD)
    ]

    coords = [
        (Inches(1.0), Inches(2.2)),
        (Inches(6.8), Inches(2.2)),
        (Inches(1.0), Inches(4.7)),
        (Inches(6.8), Inches(4.7))
    ]

    for idx, (num, title, date, tag, color) in enumerate(categories):
        left, top = coords[idx]
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.5), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = WHITE
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{num}  {title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        p2 = tf.add_paragraph()
        p2.text = f"\nDate: {date}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = YELLOW if color == RED else WHITE

        p3 = tf.add_paragraph()
        p3.text = f"Tag: {tag}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = WHITE

    # SLIDE 7: GA1 Details & Decor
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_badge(slide7, "EVENT 01")

    txBox = slide7.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "GA1 / AKWE: DETAILS & DECOR"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    card = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BLACK
    card.line.width = Pt(3)

    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "EVENT METADATA"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "• Event Head: Ash"
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Target Window: August 31 to September 4 (Poll deadline Aug 31; choices Sept 1-3)"
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Design Theme: Spider-Man"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "\nVENUE DECOR & MATERIALS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "• Banderitas: Red, blue, white, and black crepe paper theme."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Web Designs: White yarn webs across venue (reusable for Halloween)."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Accents & Photo-booth: Spider-Man printouts taped with banderitas + photo setup."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    # SLIDE 8: Interactive Bulletin Board
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_badge(slide8, "GA1 FEATURE")

    txBox = slide8.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "INTERACTIVE BULLETIN BOARD"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    pbox = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.8))
    pbox.fill.solid()
    pbox.fill.fore_color.rgb = RED
    pbox.line.color.rgb = WHITE
    pbox.line.width = Pt(2)

    tf = pbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BOARD PROMPT:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = YELLOW

    p2 = tf.add_paragraph()
    p2.text = "\"If you were to be forgotten by everyone, what will be the first thing that you will do?\""
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    dbox = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.1), Inches(11.333), Inches(2.7))
    dbox.fill.solid()
    dbox.fill.fore_color.rgb = WHITE
    dbox.line.color.rgb = BLACK
    dbox.line.width = Pt(3)

    tf = dbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EXECUTION & INCENTIVES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Materials: Masking tape, bond paper squares, markers."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Dormer Incentive: Free candies provided!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "• Process: Dormers answer the prompt on paper squares as they fall in line for GA registration."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    # SLIDE 9: GA1 Program & Party
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_badge(slide9, "GA1 PROGRAM")

    txBox = slide9.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PROGRAM & ACQUAINTANCE PARTY"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    c1 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(5.4), Inches(4.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = WHITE
    c1.line.color.rgb = BLACK
    c1.line.width = Pt(3)

    tf = c1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "GA PROGRAM FLOW"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "\n• Welcoming Remarks"
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Special Performer"
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Interactive Game (Hero/Spider-Man)"
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Presentation Deck Walkthrough"
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    c2 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.4), Inches(4.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = WHITE
    c2.line.color.rgb = BLACK
    c2.line.width = Pt(3)

    tf = c2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PARTY GAMES & PRIZES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "\nGames:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK

    games = ["Animal Sounds Team Finder", "Balloon Pop", "Charades Relay", "1-2-3 Go Game"]
    for g in games:
        p = tf.add_paragraph()
        p.text = f"  - {g}"
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "\nPrizes:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "  - Candy, Tsitsirya (snacks)"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "  - Dorm Essentials (noodles, panlaba, canned goods)"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    # SLIDE 10: Halloween Party Week
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_badge(slide10, "EVENT 02")

    txBox = slide10.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "HALLOWEEN PARTY WEEK"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    card = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BLACK
    card.line.width = Pt(3)

    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "DATE: October 19 - October 23   |   THEME: Horror"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "\nSAMPA-SPECIFIC ACTIVITIES:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK

    activities = [
        "Interactive Bulletin Board: Prompt: \"What's your horror story?\"",
        "Outdoors Movie Night: Open-air film screening",
        "Open House & Collaboration: Possible Sampa x Kamia Halloween Night collaboration",
        "Design-Your-Own-Room / Corridor Contest: Spookiest room decor competition",
        "Photobooth & Performers: Dedicated photo setup and live entertainment"
    ]

    for act in activities:
        p = tf.add_paragraph()
        p.text = f"• {act}"
        p.font.size = Pt(17)
        p.font.color.rgb = BLACK

    # SLIDE 11: GA2 / Year-End Party
    slide11 = prs.slides.add_slide(blank_layout)
    add_bg(slide11)
    add_badge(slide11, "EVENT 03")

    txBox = slide11.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "GA2 / YEAR-END PARTY"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    card = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BLACK
    card.line.width = Pt(3)

    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "DATE: End of Semester   |   THEME: Year-End Celebration"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "\nPLANNED HIGHLIGHTS:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK

    highlights = [
        "Low-Cost Exchange Gift: White Elephant style gift swap",
        "Bingo Night: Classic dorm bingo with prizes",
        "Year-End Assembly: Wrap-up celebration & recognition"
    ]

    for h in highlights:
        p = tf.add_paragraph()
        p.text = f"• {h}"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK

    # SLIDE 12: Additional Proposed Events
    slide12 = prs.slides.add_slide(blank_layout)
    add_bg(slide12)
    add_badge(slide12, "SPECIAL EVENTS")

    txBox = slide12.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "POSSIBLE ADDITIONAL EVENTS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    c1 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(5.4), Inches(4.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = WHITE
    c1.line.color.rgb = BLACK
    c1.line.width = Pt(3)

    tf = c1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "MENTAL HEALTH WEEK"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "Theme: Mental Health & The Arts"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "\n• Big canvas-style bulletin board"
    p.font.size = Pt(17)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Chalk Art activity zone"
    p.font.size = Pt(17)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Coloring pages & creative crafts"
    p.font.size = Pt(17)
    p.font.color.rgb = BLACK

    c2 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.4), Inches(4.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = WHITE
    c2.line.color.rgb = BLACK
    c2.line.width = Pt(3)

    tf = c2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "GAME NIGHTS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "Concept: Corridor Wars"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "\n• Tournament Style weekly game nights"
    p.font.size = Pt(17)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Corridor/Floor competitive format"
    p.font.size = Pt(17)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "• Cumulative prize at end of month"
    p.font.size = Pt(17)
    p.font.color.rgb = BLACK

    # SLIDE 13: Day to Day Living
    slide13 = prs.slides.add_slide(blank_layout)
    add_bg(slide13)
    add_badge(slide13, "DORM GUIDELINES")

    txBox = slide13.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "DAY TO DAY LIVING REMINDERS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    c1 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(5.4), Inches(4.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = WHITE
    c1.line.color.rgb = BLACK
    c1.line.width = Pt(3)

    tf = c1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ELECTRICAL DECLARATION"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "\nDeclare all brought electrical items:"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "  • Electric Fans & Desk Lamps"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "  • Laptops, Desktops, Printers"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "  • Cellphones & Tablets"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    c2 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.4), Inches(4.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = WHITE
    c2.line.color.rgb = BLACK
    c2.line.width = Pt(3)

    tf = c2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "STUDY & COMMON AREAS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "\n• Meals in study area discouraged; light snacks allowed."
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "\n• Observe strict cleanliness at all times."
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "\n• Be mindful of gadgets & valuable personal belongings."
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK

    # SLIDE 14: Next Steps
    slide14 = prs.slides.add_slide(blank_layout)
    add_bg(slide14)
    add_badge(slide14, "ACTION ITEMS")

    txBox = slide14.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "NEXT STEPS & POLL DEADLINE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    card = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BLACK
    card.line.width = Pt(3)

    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "IMMEDIATE ACTION ITEMS:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED

    p = tf.add_paragraph()
    p.text = "1. GA1 Date Poll: Deadline August 31 (Choices: Sept 1-3)."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "2. Material Sourcing: Crepe paper, yarn, tape, and printouts."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "3. Prize Procurement: Candy, snacks, and dorm essentials."
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

    p = tf.add_paragraph()
    p.text = "\nOPEN FLOOR FOR QUESTIONS & SUGGESTIONS!"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(base_dir, "dorm-sa-presentation-ga-1.pptx")
    prs.save(output_path)
    print(f"Successfully generated {output_path}")

if __name__ == "__main__":
    create_presentation()
